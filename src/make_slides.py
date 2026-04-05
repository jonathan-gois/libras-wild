"""
Gera apresentação PPTX dos resultados do projeto Libras 2026.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Paleta ───────────────────────────────────────────────────────────────────
DARK_BG   = RGBColor(0x1A, 0x1A, 0x2E)   # azul escuro
MID_BG    = RGBColor(0x16, 0x21, 0x3E)
ACCENT    = RGBColor(0x0F, 0x3C, 0x78)
GREEN     = RGBColor(0x2E, 0xCC, 0x71)
YELLOW    = RGBColor(0xF3, 0x9C, 0x12)
RED       = RGBColor(0xE7, 0x4C, 0x3C)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY     = RGBColor(0xCC, 0xCC, 0xCC)
DGRAY     = RGBColor(0x88, 0x88, 0x99)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]   # completamente em branco


# ── Helpers ──────────────────────────────────────────────────────────────────
def add_slide():
    sl = prs.slides.add_slide(blank)
    bg = sl.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_BG
    return sl

def txbox(sl, text, x, y, w, h,
          size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
          wrap=True, italic=False):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb

def rect(sl, x, y, w, h, fill=ACCENT, alpha=None):
    shape = sl.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def hline(sl, y, color=ACCENT):
    line = sl.shapes.add_shape(1, Inches(0.5), y, Inches(12.33), Pt(1.5))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

def add_rows(sl, data, x, y, col_widths, row_h=Inches(0.38),
             header=True, header_bg=ACCENT):
    """Tabela manual com retângulos + texto."""
    for ri, row in enumerate(data):
        bg_color = header_bg if (ri == 0 and header) else MID_BG
        if ri > 0 and ri % 2 == 0:
            bg_color = RGBColor(0x1E, 0x2A, 0x4A)
        cx = x
        for ci, (cell, cw) in enumerate(zip(row, col_widths)):
            rect(sl, cx, y, cw, row_h, fill=bg_color)
            bld  = (ri == 0 and header) or cell.startswith("**")
            txt  = cell.replace("**", "")
            col  = GREEN if "✓" in txt or "98.50" in txt or "vencedor" in txt.lower() else WHITE
            if ri == 0 and header:
                col = WHITE
            txbox(sl, txt, cx + Pt(6), y + Pt(4), cw - Pt(12), row_h - Pt(4),
                  size=11, bold=bld, color=col, align=PP_ALIGN.CENTER)
            cx += cw
        y += row_h
    return y


# ════════════════════════════════════════════════════════════════════════════
# Slide 1 — Capa
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, Inches(0), Inches(0), W, Inches(1.4), fill=RGBColor(0x0F,0x3C,0x78))
txbox(sl, "Reconhecimento de Sinais em Libras",
      Inches(0.5), Inches(0.25), Inches(12), Inches(0.9),
      size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txbox(sl, "Dataset MINDS-Libras · 800 amostras · 20 classes",
      Inches(1), Inches(1.6), Inches(11), Inches(0.5),
      size=22, color=LGRAY, align=PP_ALIGN.CENTER)

# Métricas destaque
for i, (lbl, val, col) in enumerate([
    ("MELHOR 10-FOLD CV", "98.50%", GREEN),
    ("GENERALIZAÇÃO LOPO", "74.50%", YELLOW),
    ("SUPERA ESTADO DA ARTE", "+2.5 pp", GREEN),
]):
    bx = Inches(1.0 + i * 3.8)
    rect(sl, bx, Inches(2.5), Inches(3.4), Inches(1.6),
         fill=RGBColor(0x0A,0x25,0x55))
    txbox(sl, val, bx, Inches(2.6), Inches(3.4), Inches(0.8),
          size=40, bold=True, color=col, align=PP_ALIGN.CENTER)
    txbox(sl, lbl, bx, Inches(3.3), Inches(3.4), Inches(0.6),
          size=11, color=DGRAY, align=PP_ALIGN.CENTER)

txbox(sl, "ExtraTrees-1000  ·  GEI + FFT Sliding Window + Kinematic Stats  ·  11.172 features",
      Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.5),
      size=16, color=LGRAY, align=PP_ALIGN.CENTER, italic=True)

txbox(sl, "Hardware: CPU  ·  Sem GPU  ·  Sem vídeo RGB  ·  Apenas skeleton MediaPipe",
      Inches(0.5), Inches(5.1), Inches(12.3), Inches(0.4),
      size=13, color=DGRAY, align=PP_ALIGN.CENTER)

txbox(sl, "Projeto Libras 2026", Inches(0.5), Inches(6.9), Inches(6), Inches(0.4),
      size=11, color=DGRAY)
txbox(sl, "2026", Inches(11.5), Inches(6.9), Inches(1.3), Inches(0.4),
      size=11, color=DGRAY, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════════
# Slide 2 — Contexto & Dataset
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, Inches(0), Inches(0), W, Inches(0.9), fill=ACCENT)
txbox(sl, "Contexto & Dataset", Inches(0.4), Inches(0.1), Inches(10), Inches(0.7),
      size=28, bold=True)
hline(sl, Inches(0.95))

col1x, col2x = Inches(0.5), Inches(6.8)
cw = Inches(5.9)

txbox(sl, "MINDS-Libras", col1x, Inches(1.1), cw, Inches(0.4),
      size=18, bold=True, color=GREEN)
items1 = [
    "800 amostras  (original) / 1180 com Zenodo",
    "20 sinais isolados de Libras",
    "10 sinalizadores × 40 amostras cada",
    "MediaPipe Holistic: 75 landmarks × 3 coords",
    "= 225 features por frame",
]
for i, t in enumerate(items1):
    txbox(sl, "▸  " + t, col1x + Inches(0.2), Inches(1.55 + i*0.42), cw, Inches(0.4),
          size=14, color=LGRAY)

txbox(sl, "Classes (20 sinais)", col1x, Inches(3.65), cw, Inches(0.4),
      size=15, bold=True, color=YELLOW)
classes = ["acontecer","aluno","amarelo","america","aproveitar",
           "bala","banco","banheiro","barulho","cinco",
           "conhecer","espelho","esquina","filho","maca",
           "medo","ruim","sapo","vacina","vontade"]
cols_c = 4
for i, c in enumerate(classes):
    bx = col1x + Inches((i % cols_c) * 1.42)
    by = Inches(4.1 + (i // cols_c) * 0.38)
    rect(sl, bx, by, Inches(1.35), Inches(0.33), fill=RGBColor(0x0F,0x2D,0x55))
    txbox(sl, c, bx + Pt(4), by + Pt(2), Inches(1.35), Inches(0.33),
          size=10, color=LGRAY, align=PP_ALIGN.CENTER)

txbox(sl, "Pipeline de Extração", col2x, Inches(1.1), cw, Inches(0.4),
      size=18, bold=True, color=GREEN)
steps = [
    ("1", "Vídeo MP4", "MediaPipe Holistic extrai 75 landmarks/frame"),
    ("2", "Normalização", "Centro: quadris  |  Escala: distância ombro-ombro"),
    ("3", "GEI 64×48px", "Esqueleto 2D médio por pixel → 3.072 dims"),
    ("4", "Kinematic Stats", "mean/std/max/min de vel e acel → 1.800 dims"),
    ("5", "FFT Sliding", "janela=16, passo=8, top-5 mag → 6.300 dims"),
    ("6", "Concat", "Total: 11.172 dims  →  ExtraTrees-1000"),
]
for i, (num, title, desc) in enumerate(steps):
    by = Inches(1.55 + i * 0.82)
    rect(sl, col2x, by, Inches(0.38), Inches(0.38),
         fill=RGBColor(0x0F,0x3C,0x78))
    txbox(sl, num, col2x, by, Inches(0.38), Inches(0.38),
          size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(sl, title, col2x + Inches(0.44), by, Inches(1.5), Inches(0.38),
          size=13, bold=True, color=WHITE)
    txbox(sl, desc, col2x + Inches(0.44), by + Inches(0.28), Inches(5.2), Inches(0.4),
          size=11, color=DGRAY)


# ════════════════════════════════════════════════════════════════════════════
# Slide 3 — Comparativo Geral
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, Inches(0), Inches(0), W, Inches(0.9), fill=ACCENT)
txbox(sl, "Comparativo — Estado da Arte vs Nossos Modelos",
      Inches(0.4), Inches(0.1), Inches(12), Inches(0.7), size=26, bold=True)
hline(sl, Inches(0.95))

data = [
    ["Método", "Acurácia", "Δ vs SotA (96%)"],
    ["Passos et al. — GEI + SVM", "84.66%", "−11.3 pp"],
    ["Rego et al. — FFT+Kin+BiLSTM", "92.0%", "−4.0 pp"],
    ["Rezende et al. — CNN 3D", "93.3%", "−2.7 pp"],
    ["De Castro / Arcanjo — 3D CNN / FastDTW", "~96.0%", "referência"],
    ["── Nossos modelos ──────────────────", "", ""],
    ["BiLSTM multimodal (sem augmentação)", "72.50% ±5.51%", "−23.5 pp"],
    ["Transformer + mirror aug (inválido)", "88.75% ±4.51%", "−7.3 pp"],
    ["ST-GCN + MobileNetV3 (v5)", "91.88% ±4.12%", "−4.1 pp"],
    ["Transformer sem mirror (v4)", "94.62% ±5.59%", "−1.4 pp"],
    ["ExtraTrees-1000 (1180 amostras)", "96.78% ±2.03%", "+0.8 pp"],
    ["**✓  ExtraTrees-1000 (800 amostras)", "**98.50% ±0.75%", "**+2.5 pp"],
]
cws = [Inches(6.5), Inches(2.8), Inches(2.8)]
add_rows(sl, data, Inches(0.5), Inches(1.05), cws,
         row_h=Inches(0.44), header=True)

txbox(sl, "⚠  LOPO (generalização real): 74.50% ±11.68%  —  gap de 24 pp",
      Inches(0.5), Inches(6.85), Inches(12), Inches(0.45),
      size=13, color=YELLOW, bold=True, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# Slide 4 — Método Vencedor
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, Inches(0), Inches(0), W, Inches(0.9), fill=GREEN)
txbox(sl, "Método Vencedor — ExtraTrees-1000",
      Inches(0.4), Inches(0.1), Inches(12), Inches(0.7),
      size=28, bold=True, color=DARK_BG)
hline(sl, Inches(0.95), color=GREEN)

# Parâmetros
txbox(sl, "Parâmetros", Inches(0.5), Inches(1.05), Inches(5.8), Inches(0.4),
      size=16, bold=True, color=GREEN)
params = [
    ["Parâmetro", "Valor"],
    ["n_estimators", "1000"],
    ["criterion", "gini"],
    ["max_features", "sqrt"],
    ["min_samples_leaf", "1"],
    ["bootstrap", "False  (diferença chave vs RF)"],
    ["n_jobs", "−1  (todos os núcleos)"],
    ["random_state", "42"],
]
add_rows(sl, params, Inches(0.5), Inches(1.5), [Inches(2.5), Inches(3.2)],
         row_h=Inches(0.4))

# Por que vence
txbox(sl, "Por que supera redes neurais?", Inches(6.8), Inches(1.05),
      Inches(6.0), Inches(0.4), size=16, bold=True, color=GREEN)
reasons = [
    "800 amostras / 20 classes = apenas 40 por classe",
    "Redes neurais: ~9M parâmetros → overfitting severo",
    "ExtraTrees: ~100K parâmetros efetivos",
    "Sem PCA: lida nativamente com 11.172 dims",
    "Thresholds aleatórios → menor variância que RF",
    "Treinamento em segundos (CPU)",
]
for i, t in enumerate(reasons):
    txbox(sl, "▸  " + t, Inches(7.0), Inches(1.5 + i * 0.42),
          Inches(5.8), Inches(0.4), size=13, color=LGRAY)

# Importância das features
hline(sl, Inches(4.15), color=DGRAY)
txbox(sl, "Importância das Features (Gini)", Inches(0.5), Inches(4.25),
      Inches(12), Inches(0.4), size=16, bold=True, color=YELLOW)

bars = [
    ("FFT Sliding Window  (6.300 dims)", 74.95, GREEN),
    ("Kinematic Stats  (1.800 dims)",    13.32, YELLOW),
    ("GEI  (3.072 dims)",                11.74, RGBColor(0xE7,0x4C,0x3C)),
]
bar_x = Inches(0.5)
for i, (lbl, pct, col) in enumerate(bars):
    by = Inches(4.75 + i * 0.7)
    txbox(sl, lbl, bar_x, by, Inches(3.0), Inches(0.35), size=12, color=LGRAY)
    bw = Inches(0.085 * pct)
    rect(sl, bar_x + Inches(3.1), by + Pt(2), bw, Inches(0.3), fill=col)
    txbox(sl, f"{pct:.1f}%", bar_x + Inches(3.1) + bw + Pt(6), by,
          Inches(0.7), Inches(0.35), size=13, bold=True, color=col)

txbox(sl, "FFT sozinho já atinge 98.12%  —  combinação reduz variância ±1.28% → ±0.75%",
      Inches(0.5), Inches(7.0), Inches(12), Inches(0.38),
      size=12, color=DGRAY, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# Slide 5 — Mirror Augmentation: Achado Crítico
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, Inches(0), Inches(0), W, Inches(0.9), fill=RED)
txbox(sl, "Achado Crítico — Mirror Augmentation é Prejudicial",
      Inches(0.4), Inches(0.1), Inches(12), Inches(0.7),
      size=26, bold=True, color=WHITE)
hline(sl, Inches(0.95), color=RED)

txbox(sl, "Hipótese inicial:", Inches(0.5), Inches(1.1), Inches(12), Inches(0.4),
      size=15, bold=True, color=LGRAY)
txbox(sl, '"Espelhar o sinal simula sinalizadores canhotos → dobra o dataset → melhora generalização"',
      Inches(0.5), Inches(1.5), Inches(12), Inches(0.5),
      size=14, color=LGRAY, italic=True)

hline(sl, Inches(2.1), color=DGRAY)
txbox(sl, "Resultado da análise k-NN no espaço FFT:", Inches(0.5), Inches(2.2),
      Inches(12), Inches(0.4), size=15, bold=True, color=YELLOW)

findings = [
    ("94.9%", "das amostras espelhadas têm vizinho mais próximo em OUTRA classe"),
    ("19 / 20", "classes têm seu espelho mapeado majoritariamente para outra classe"),
    ("Única exceção:", '"acontecer" — sinal simétrico (palmas frente a frente)'),
]
for i, (num, desc) in enumerate(findings):
    by = Inches(2.7 + i * 0.62)
    rect(sl, Inches(0.5), by, Inches(1.5), Inches(0.48), fill=RGBColor(0x4A,0x10,0x10))
    txbox(sl, num, Inches(0.5), by, Inches(1.5), Inches(0.48),
          size=18, bold=True, color=RED, align=PP_ALIGN.CENTER)
    txbox(sl, desc, Inches(2.1), by + Pt(6), Inches(10.5), Inches(0.4),
          size=14, color=LGRAY)

hline(sl, Inches(4.65), color=DGRAY)
txbox(sl, "Causa fonológica:", Inches(0.5), Inches(4.75),
      Inches(12), Inches(0.4), size=15, bold=True, color=YELLOW)
txbox(sl, "Libras usa lateralidade como traço fonológico. Espelhar não cria variante válida — cria confusão com sinais opostos.",
      Inches(0.5), Inches(5.2), Inches(12), Inches(0.5), size=14, color=LGRAY)

# Impacto
data = [
    ["Modelo", "Augmentação", "Acurácia", "Δ"],
    ["v3 Transformer", "COM mirror", "88.75% ±4.51%", "baseline"],
    ["v4 Transformer", "SEM mirror", "94.62% ±5.59%", "+5.87 pp  ✓"],
]
add_rows(sl, data, Inches(2.0), Inches(5.85),
         [Inches(3.0), Inches(2.5), Inches(2.8), Inches(2.5)],
         row_h=Inches(0.45))


# ════════════════════════════════════════════════════════════════════════════
# Slide 6 — Gap de Generalização (LOPO)
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, Inches(0), Inches(0), W, Inches(0.9), fill=YELLOW)
txbox(sl, "Generalização Real — Leave-One-Person-Out (LOPO)",
      Inches(0.4), Inches(0.1), Inches(12), Inches(0.7),
      size=26, bold=True, color=DARK_BG)
hline(sl, Inches(0.95), color=YELLOW)

# Gap visual
rect(sl, Inches(0.5), Inches(1.1), Inches(5.8), Inches(2.0),
     fill=RGBColor(0x07,0x2A,0x0F))
txbox(sl, "10-fold CV", Inches(0.5), Inches(1.15), Inches(5.8), Inches(0.5),
      size=14, color=LGRAY, align=PP_ALIGN.CENTER)
txbox(sl, "98.50%", Inches(0.5), Inches(1.65), Inches(5.8), Inches(0.9),
      size=52, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
txbox(sl, "mesmo sinalizador no treino e no teste", Inches(0.5), Inches(2.85),
      Inches(5.8), Inches(0.3), size=11, color=DGRAY, align=PP_ALIGN.CENTER)

rect(sl, Inches(7.0), Inches(1.1), Inches(5.8), Inches(2.0),
     fill=RGBColor(0x3A,0x1A,0x05))
txbox(sl, "LOPO", Inches(7.0), Inches(1.15), Inches(5.8), Inches(0.5),
      size=14, color=LGRAY, align=PP_ALIGN.CENTER)
txbox(sl, "74.50%", Inches(7.0), Inches(1.65), Inches(5.8), Inches(0.9),
      size=52, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
txbox(sl, "sinalizador nunca visto no treino", Inches(7.0), Inches(2.85),
      Inches(5.8), Inches(0.3), size=11, color=DGRAY, align=PP_ALIGN.CENTER)

txbox(sl, "↓ 24 pp", Inches(6.0), Inches(1.85), Inches(1.3), Inches(0.6),
      size=20, bold=True, color=RED, align=PP_ALIGN.CENTER)

# Por sinalizador
txbox(sl, "Resultado por sinalizador (ET-1000)", Inches(0.5), Inches(3.3),
      Inches(12), Inches(0.4), size=15, bold=True, color=YELLOW)
sdata = [
    ["Sinalizador", "Acurácia"],
    ["Sinalizador01", "77.5%"],
    ["Sinalizador02", "88.0%"],
    ["Sinalizador05", "76.5%"],
    ["Sinalizador06", "82.5%"],
    ["Sinalizador08", "75.0%"],
    ["Sinalizador10", "51.0%  ⚠"],
    ["Média", "74.50% ±11.68%"],
]
add_rows(sl, sdata, Inches(0.5), Inches(3.75),
         [Inches(3.5), Inches(2.5)], row_h=Inches(0.42))

txbox(sl, "Interpretação:", Inches(6.8), Inches(3.35), Inches(6.0), Inches(0.4),
      size=14, bold=True, color=YELLOW)
interp = [
    "O modelo aprende o estilo motor de cada sinalizador,",
    "não as propriedades fonológicas universais do sinal.",
    "",
    "Soluções:",
    "▸  Mais sinalizadores (Zenodo: +4 → 1180 amostras)",
    "▸  Normalização adversarial por sinalizador",
    "▸  Features de velocidade relativa (invariância ao estilo)",
    "▸  Data augmentation estilo-agnóstica",
]
for i, t in enumerate(interp):
    col = LGRAY if not t.startswith("▸") else WHITE
    txbox(sl, t, Inches(6.8), Inches(3.78 + i * 0.38),
          Inches(6.0), Inches(0.36), size=12, color=col)


# ════════════════════════════════════════════════════════════════════════════
# Slide 7 — Ablação & Pipeline YouTube
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, Inches(0), Inches(0), W, Inches(0.9), fill=ACCENT)
txbox(sl, "Ablação de Features  ·  Pipeline YouTube",
      Inches(0.4), Inches(0.1), Inches(12), Inches(0.7), size=26, bold=True)
hline(sl, Inches(0.95))

txbox(sl, "Contribuição por Modalidade (ET-1000)", Inches(0.5), Inches(1.05),
      Inches(6.0), Inches(0.4), size=15, bold=True, color=GREEN)
abl = [
    ["Features", "Acurácia", "±std"],
    ["GEI only", "52.00%", "±4.00%"],
    ["KIN only", "81.62%", "±2.44%"],
    ["FFT only", "98.12%", "±1.28%"],
    ["GEI + KIN", "81.62%", "±5.39%"],
    ["GEI + FFT", "98.00%", "±1.50%"],
    ["KIN + FFT", "98.00%", "±1.27%"],
    ["**✓  GEI + KIN + FFT", "**98.50%", "**±0.75%"],
]
add_rows(sl, abl, Inches(0.5), Inches(1.5),
         [Inches(3.0), Inches(1.6), Inches(1.2)], row_h=Inches(0.42))

txbox(sl, "FFT captura ~75% da informação.\nCombinação completa reduz variância em 40%.",
      Inches(0.5), Inches(5.15), Inches(6.0), Inches(0.6),
      size=12, color=DGRAY, italic=True)

# YouTube
txbox(sl, "Pipeline YouTube — Análise de Vídeo em Libras",
      Inches(6.9), Inches(1.05), Inches(6.0), Inches(0.4),
      size=15, bold=True, color=GREEN)
yt_steps = [
    ("yt-dlp", "Download URL → MP4 + SRT legendas"),
    ("MediaPipe", "Extração de landmarks (30.775 frames, 17 min)"),
    ("Sign Spotter", "Energia vel. punho + threshold → 559 segmentos"),
    ("ET-1000", "Classificação top-5 por segmento"),
    ("Alinhamento", "Overlap ≥30% com legendas → 437 entradas"),
]
for i, (comp, desc) in enumerate(yt_steps):
    by = Inches(1.5 + i * 0.62)
    rect(sl, Inches(6.9), by, Inches(1.5), Inches(0.48),
         fill=RGBColor(0x0F,0x3C,0x78))
    txbox(sl, comp, Inches(6.9), by, Inches(1.5), Inches(0.48),
          size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(sl, desc, Inches(8.5), by + Pt(6), Inches(4.5), Inches(0.4),
          size=12, color=LGRAY)

txbox(sl, "Top sinais detectados:", Inches(6.9), Inches(4.65),
      Inches(5.8), Inches(0.4), size=13, bold=True, color=YELLOW)
yt_res = [
    ("vacina", "35.8%"), ("acontecer", "15.1%"), ("medo", "14.0%"),
    ("aproveitar", "13.8%"), ("esquina", "13.3%"),
]
for i, (sign, pct) in enumerate(yt_res):
    bx = Inches(6.9 + (i % 3) * 2.1)
    by = Inches(5.05 + (i // 3) * 0.52)
    rect(sl, bx, by, Inches(1.95), Inches(0.42), fill=RGBColor(0x0F,0x2D,0x55))
    txbox(sl, f"{sign}  {pct}", bx + Pt(4), by + Pt(4), Inches(1.9), Inches(0.38),
          size=12, color=WHITE, align=PP_ALIGN.CENTER)

txbox(sl, "Limitação: 20 sinais conhecidos. Confianças baixas (13–21%) indicam\n"
          "maioria dos sinais fora do vocabulário do modelo.",
      Inches(6.9), Inches(5.95), Inches(6.0), Inches(0.7),
      size=11, color=DGRAY, italic=True)


# ════════════════════════════════════════════════════════════════════════════
# Slide 8 — Conclusões & Próximos Passos
# ════════════════════════════════════════════════════════════════════════════
sl = add_slide()
rect(sl, Inches(0), Inches(0), W, Inches(0.9), fill=ACCENT)
txbox(sl, "Conclusões & Próximos Passos",
      Inches(0.4), Inches(0.1), Inches(12), Inches(0.7), size=28, bold=True)
hline(sl, Inches(0.95))

txbox(sl, "Conclusões", Inches(0.5), Inches(1.1), Inches(6.0), Inches(0.4),
      size=16, bold=True, color=GREEN)
concl = [
    "ExtraTrees-1000 supera SotA (+2.5pp) sem GPU, sem RGB",
    "FFT domina: 74.95% da importância, 98.12% sozinho",
    "Mirror augmentation destrói performance (−5.87pp)",
    "Gap CV vs LOPO (24pp): modelo memoriza sinalizador",
    "ST-GCN não supera Transformer simples com 800 amostras",
    "Novos sinalizadores tornam o CV mais realista (96.78%)",
]
for i, t in enumerate(concl):
    txbox(sl, "✓  " + t if i < 3 else "⚠  " + t,
          Inches(0.5), Inches(1.55 + i * 0.48), Inches(6.0), Inches(0.45),
          size=13, color=GREEN if i < 3 else YELLOW)

txbox(sl, "Próximos Passos", Inches(7.0), Inches(1.1), Inches(6.0), Inches(0.4),
      size=16, bold=True, color=YELLOW)
next_steps = [
    "LOPO com 1180 amostras (10 sinalizadores completos)",
    "Anotação fonológica com Qwen2-VL-7B (CM, PA, MOV, OR, ENM)",
    "v6: PhoneticFusionModel + InfoNCE contrastive loss",
    "Normalização adversarial por sinalizador",
    "Threshold de rejeição (out-of-vocabulary) no pipeline YouTube",
    "Expansão do vocabulário além das 20 classes",
]
for i, t in enumerate(next_steps):
    txbox(sl, "→  " + t, Inches(7.0), Inches(1.55 + i * 0.48),
          Inches(6.0), Inches(0.45), size=13, color=LGRAY)

# Resumo final
hline(sl, Inches(4.6), color=GREEN)
rect(sl, Inches(0.5), Inches(4.75), Inches(12.3), Inches(1.0),
     fill=RGBColor(0x07,0x2A,0x0F))
txbox(sl,
      "Com apenas features clássicas (FFT + kinematic stats + GEI) e ExtraTrees-1000, "
      "atingimos 98.50% ±0.75% em 10-fold CV — superando todo o estado da arte em +2.5pp "
      "sem usar vídeo RGB, sem GPU e sem redes profundas.\n"
      "O gargalo real é a generalização entre sinalizadores (74.50% LOPO).",
      Inches(0.7), Inches(4.85), Inches(11.9), Inches(0.85),
      size=14, color=WHITE, italic=True)

# Rodapé
txbox(sl, "Projeto Libras 2026  ·  MINDS-Libras  ·  Resultados acumulados",
      Inches(0.5), Inches(6.9), Inches(12), Inches(0.4),
      size=11, color=DGRAY, align=PP_ALIGN.CENTER)


# ── Salva ────────────────────────────────────────────────────────────────────
out = "results/libras2026_resultados.pptx"
prs.save(out)
print(f"Salvo: {out}")
